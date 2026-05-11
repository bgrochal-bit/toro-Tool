import streamlit as st
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont
import datetime
import io
import os
import base64
import requests
from openai import OpenAI
import json

# --- INITIALIZE OpenAI ---
client = OpenAI(api_key=st.secrets["OPENAI_API_KEY"]) if "OPENAI_API_KEY" in st.secrets else None

# --- TORO BRANDING ---
TORO_NAVY = (0, 49, 82)      # RGB for PDF/PPTX
TORO_NAVY_HEX = "#003152"    # Hex for CSS
TORO_GOLD = (244, 194, 68)   # RGB for PDF/PPTX
TORO_GOLD_HEX = "#f4c244"    # Hex for CSS

st.set_page_config(page_title="Toro AI Design Engine", layout="wide")

# --- CUSTOM CSS ---
st.markdown(f"""
    <style>
    .stApp {{ background-color: {TORO_NAVY_HEX}; color: white; }}
    .stButton>button {{ background-color: {TORO_GOLD_HEX} !important; color: {TORO_NAVY_HEX} !important; font-weight: bold !important; width: 100%; border-radius: 4px; }}
    .stTextArea>div>textarea, .stTextInput>div>div>input {{ background-color: rgba(255,255,255,0.1) !important; color: white !important; border: 1px solid rgba(255,255,255,0.2) !important; }}
    h1, h2, h3, p, label {{ color: white !important; font-family: 'Inter', sans-serif; }}
    </style>
    """, unsafe_allow_html=True)

# --- HELPERS ---
def encode_image(image_bytes):
    return base64.b64encode(image_bytes).decode('utf-8')

def download_image(url):
    response = requests.get(url)
    return io.BytesIO(response.content)

# --- ENGINE: SYNTHESIS (DALL-E 3 + GPT-4o Vision) ---
def synthesize_product(uploaded_files, instructions):
    if not client: return None
    # Vision Architect: GPT-4o analyzes uploads to create a perfect prompt for DALL-E 3
    msgs = [{"role": "user", "content": [{"type": "text", "text": f"Analyze these references. Generate a detailed prompt for DALL-E 3 to synthesize ONE new product image based on these instructions: {instructions}. Be technical and specific."}]}]
    
    for f in uploaded_files:
        msgs[0]["content"].append({"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{encode_image(f.getvalue())}"}})
    
    analysis = client.chat.completions.create(model="gpt-4o", messages=msgs)
    prompt = analysis.choices[0].message.content
    
    gen = client.images.generate(model="dall-e-3", prompt=prompt, size="1024x1024", quality="hd")
    return gen.data[0].url

# --- ENGINE: SMART MARKUP ARROWS ---
def generate_markup_image(base_img_bytes, feedback_text):
    if not client: return None, []
    
    prompt = f"""
    Look at this product image. Based on these comments: {feedback_text}
    Identify the X and Y coordinates (0-1000 scale) for where an arrow should point for each comment.
    Return ONLY a JSON object with a key 'points' containing a list of: {{"id": 1, "x": 500, "y": 300, "text": "comment text"}}
    """
    
    response = client.chat.completions.create(
        model="gpt-4o",
        messages=[{"role": "user", "content": [{"type": "text", "text": prompt}, {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{encode_image(base_img_bytes)}"}}]}],
        response_format={ "type": "json_object" }
    )
    
    data = json.loads(response.choices[0].message.content)
    points = data.get("points", [])
    
    # Draw on image
    img = Image.open(io.BytesIO(base_img_bytes)).convert("RGB")
    draw = ImageDraw.Draw(img)
    w, h = img.size
    
    for p in points:
        nx, ny = (p['x']/1000)*w, (p['y']/1000)*h
        # Draw Arrow (Toro Gold)
        draw.line([(nx - 60, ny - 60), (nx, ny)], fill=(244, 194, 68), width=10)
        # Draw Label Circle
        draw.ellipse([nx-80, ny-80, nx-40, ny-40], fill=(0, 49, 82), outline=(244, 194, 68), width=3)
        # In a real app, you'd add text here; for now, we leave the label circle
        
    out_io = io.BytesIO()
    img.save(out_io, format="PNG")
    return out_io.getvalue(), points

# --- EXPORT ENGINE (PDF & PPTX) ---
def export_docs(prod_name, version, mode, img_bytes, comments):
    # PPTX
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Navy Header
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.8))
    rect.fill.solid()
    rect.fill.foreground_color.rgb = RGBColor(*TORO_NAVY)
    
    title = slide.shapes.add_textbox(Inches(0.2), Inches(0.1), Inches(8), Inches(0.5))
    title.text_frame.text = f"TORO | {mode.upper()} | {version}"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255,255,255)
    title.text_frame.paragraphs[0].font.bold = True
    
    # Image
    slide.shapes.add_picture(io.BytesIO(img_bytes), Inches(0.5), Inches(1.2), height=Inches(5))
    
    # Comments (Editable Shapes)
    if mode == "Feedback Card Generator":
        for i, c in enumerate(comments):
            # Editable Text Box
            tb = slide.shapes.add_textbox(Inches(6.5), Inches(1.5 + (i*0.8)), Inches(3), Inches(0.7))
            tb.text_frame.text = f"{c['id']}. {c['text']}"
            tb.text_frame.paragraphs[0].font.size = Pt(12)
            # Editable Arrow
            arrow = slide.shapes.add_shape(MSO_SHAPE.BENT_ARROW, Inches(6), Inches(1.6 + (i*0.8)), Inches(0.4), Inches(0.4))
            arrow.fill.solid()
            arrow.fill.foreground_color.rgb = RGBColor(*TORO_GOLD)

    ppt_io = io.BytesIO()
    prs.save(ppt_io)
    return ppt_io.getvalue()

# --- APP UI ---
st.title("Simplify Your Design. Scale Your Business.")
mode = st.radio("Choose Workflow", ["Tech Pack Generator", "Feedback Card Generator"], horizontal=True)

col1, col2 = st.columns([1, 1.2])

with col1:
    st.subheader("1. Design Inputs")
    prod_name = st.text_input("Product Name", value="")
    
    # TP = P1 always, FB = Choice
    version = "P1" if mode == "Tech Pack Generator" else st.selectbox("Select Feedback Round", ["R1", "R2", "R3"])
    
    uploads = st.file_uploader("Upload Image References", accept_multiple_files=True)
    
    if mode == "Tech Pack Generator":
        notes = st.text_area("Synthesis Instructions", placeholder="e.g. Combine the shape of image 1 with the color of image 2...")
        if st.button("SYNTHESIZE PRODUCT"):
            url = synthesize_product(uploads, notes)
            st.session_state['gen_bytes'] = download_image(url).getvalue()
    else:
        feedback = st.text_area("Feedback Details", placeholder="1. Move logo higher\n2. Change button color...")
        if st.button("GENERATE FEEDBACK MARKUP"):
            res_bytes, pts = generate_markup_image(uploads[0].getvalue(), feedback)
            st.session_state['fb_bytes'] = res_bytes
            st.session_state['fb_points'] = pts

with col2:
    st.subheader("2. Visual Output")
    active_img_bytes = None
    
    if mode == "Tech Pack Generator" and 'gen_bytes' in st.session_state:
        active_img_bytes = st.session_state['gen_bytes']
        st.image(active_img_bytes, caption="AI Synthesized Result (DALL-E 3)")
    elif mode == "Feedback Card Generator" and 'fb_bytes' in st.session_state:
        active_img_bytes = st.session_state['fb_bytes']
        st.image(active_img_bytes, caption="AI Markup Result (GPT-4o Vision)")
    elif uploads:
        active_img_bytes = uploads[0].getvalue()
        st.image(active_img_bytes, caption="Reference Upload")

    if active_img_bytes:
        st.write("---")
        pptx = export_docs(prod_name, version, mode, active_img_bytes, st.session_state.get('fb_points', []))
        st.download_button("📥 DOWNLOAD EDITABLE PPTX", data=pptx, file_name=f"Toro_{prod_name}.pptx")
